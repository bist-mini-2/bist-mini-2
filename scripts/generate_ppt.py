import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR

# =============================================================================
# Pastel Sage Color Palette (mirrors frontend theme.css)
# =============================================================================
BG_COLOR        = RGBColor(252, 251, 246)  # #FCFBF6  Warm Ivory
FG_COLOR        = RGBColor( 46,  51,  46)  # #2E332E  Dark Charcoal
ACCENT_COLOR    = RGBColor(142, 157, 135)  # #8E9D87  Sage Green
TERTIARY_COLOR  = RGBColor(213, 220, 208)  # #D5DCD0  Light Sage
WHITE_COLOR     = RGBColor(255, 255, 255)  # #FFFFFF  Card White
BORDER_COLOR    = RGBColor(235, 231, 224)  # #EBE7E0  Card Border
MUTED_COLOR     = RGBColor(115, 120, 115)  # Soft gray body text
SIDEBAR_COLOR   = RGBColor(250, 246, 240)  # #FAF6F0  Sidebar tone
DEEP_ACCENT     = RGBColor(100, 112,  94)  # Darker sage for variety

# =============================================================================
# Helpers
# =============================================================================
def init_presentation():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def bg_rect(prs, slide, color):
    """Full-slide background fill."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    return r

def rect(slide, x, y, w, h, fill, line_color=None, line_width=Pt(0)):
    """Generic rectangle helper."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def rrect(slide, x, y, w, h, fill=WHITE_COLOR, border=BORDER_COLOR, bw=Pt(1.2)):
    """Rounded rectangle card."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border; shape.line.width = bw
    try: shape.adjustments[0] = 0.03
    except: pass
    return shape

def txbox(slide, x, y, w, h,
          text="", font_size=Pt(9.5), bold=False, italic=False,
          color=FG_COLOR, align=PP_ALIGN.LEFT,
          v_anchor=MSO_ANCHOR.TOP,
          font_name="Malgun Gothic",
          line_spacing=1.25,
          margin_l=Inches(0.2), margin_r=Inches(0.2),
          margin_t=Inches(0.12), margin_b=Inches(0.12),
          wrap=True):
    """
    Single-paragraph text box. Supports vertical anchoring.
    For multi-paragraph use txbox_para().
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap        = wrap
    tf.auto_size        = None
    tf.vertical_anchor  = v_anchor
    tf.margin_left      = margin_l
    tf.margin_right     = margin_r
    tf.margin_top       = margin_t
    tf.margin_bottom    = margin_b

    p = tf.paragraphs[0]
    p.text          = text
    p.font.size     = font_size
    p.font.bold     = bold
    p.font.italic   = italic
    p.font.color.rgb = color
    p.font.name     = font_name
    p.alignment     = align
    p.line_spacing  = line_spacing
    return tb

def txbox_para(slide, x, y, w, h,
               paragraphs,        # list of dicts: {text, size, bold, color, space_before, align, italic, indent}
               v_anchor=MSO_ANCHOR.TOP,
               margin_l=Inches(0.2), margin_r=Inches(0.2),
               margin_t=Inches(0.15), margin_b=Inches(0.15),
               default_font="Malgun Gothic",
               default_ls=1.3):
    """
    Multi-paragraph text box with per-paragraph formatting.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap        = True
    tf.auto_size        = None
    tf.vertical_anchor  = v_anchor
    tf.margin_left      = margin_l
    tf.margin_right     = margin_r
    tf.margin_top       = margin_t
    tf.margin_bottom    = margin_b

    for i, pg in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text          = pg.get("text", "")
        p.font.size     = pg.get("size", Pt(9.5))
        p.font.bold     = pg.get("bold", False)
        p.font.italic   = pg.get("italic", False)
        p.font.color.rgb = pg.get("color", FG_COLOR)
        p.font.name     = pg.get("font", default_font)
        p.alignment     = pg.get("align", PP_ALIGN.LEFT)
        p.line_spacing  = pg.get("ls", default_ls)
        p.space_before  = pg.get("before", Pt(0))
        p.space_after   = pg.get("after", Pt(0))
    return tb

def slide_footer(slide, page_num, total=8):
    """Bottom right page indicator footer."""
    rect(slide, Inches(11.8), Inches(7.15), Inches(1.3), Inches(0.25), BG_COLOR)
    txbox(slide, Inches(11.8), Inches(7.15), Inches(1.3), Inches(0.25),
          text=f"{page_num} / {total}", font_size=Pt(8.5),
          color=MUTED_COLOR, align=PP_ALIGN.RIGHT)

def card_accent_bar(slide, x, y, w, color=ACCENT_COLOR, h=Inches(0.045)):
    """Thin top accent bar on cards."""
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()

def section_label(slide, x, y, w, label, label_color=ACCENT_COLOR):
    """Uppercase section category label above slide title."""
    txbox(slide, x, y, w, Inches(0.28),
          text=label.upper(), font_size=Pt(8.5), bold=True,
          color=label_color, letter_spacing_workaround=None)

def slide_header(slide, category_text, title_text, cat_x=Inches(0.8), title_x=Inches(0.8)):
    """Consistent slide header: category + horizontal rule + big title."""
    # Category label
    txbox(slide, cat_x, Inches(0.32), Inches(11.7), Inches(0.28),
          text=category_text.upper(), font_size=Pt(8.5), bold=True, color=ACCENT_COLOR)
    # Full-width thin rule under category
    rect(slide, Inches(0.8), Inches(0.60), Inches(11.73), Inches(0.025), TERTIARY_COLOR)
    # Main title
    txbox(slide, title_x, Inches(0.65), Inches(11.7), Inches(0.75),
          text=title_text, font_size=Pt(26), bold=True, color=FG_COLOR,
          v_anchor=MSO_ANCHOR.MIDDLE)

# =============================================================================
# Slide 1: Cover
# =============================================================================
def slide_1_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rect(prs, slide, BG_COLOR)

    # Left sage panel (1/3 width)
    panel_w = Inches(4.2)
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, panel_w, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = ACCENT_COLOR
    r.line.fill.background()

    # Subtle lighter diagonal accent inside left panel
    dec = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.4), Inches(0.4),
                                 Inches(3.2), Inches(6.7))
    dec.fill.solid(); dec.fill.fore_color.rgb = DEEP_ACCENT
    dec.line.fill.background()
    try: dec.adjustments[0] = 0.04
    except: pass

    # Vertical label on left panel
    txbox(slide, Inches(0.45), Inches(0.7), Inches(3.5), Inches(1.1),
          text="BIST MINI-2 PROJECT", font_size=Pt(9.5), bold=True,
          color=WHITE_COLOR, align=PP_ALIGN.LEFT,
          v_anchor=MSO_ANCHOR.MIDDLE)

    # Project title (left panel, Korean)
    txbox_para(slide,
               Inches(0.45), Inches(1.7), Inches(3.6), Inches(3.0),
               [
                   {"text": "웹 서비스 및\nRAG 기반\nAI 솔루션",
                    "size": Pt(30), "bold": True, "color": WHITE_COLOR,
                    "ls": 1.4, "align": PP_ALIGN.LEFT},
               ],
               v_anchor=MSO_ANCHOR.MIDDLE)

    # Bottom date tag
    txbox(slide, Inches(0.45), Inches(6.5), Inches(3.4), Inches(0.5),
          text="2026. 06. 30", font_size=Pt(11),
          color=TERTIARY_COLOR, align=PP_ALIGN.LEFT,
          v_anchor=MSO_ANCHOR.MIDDLE)

    # Right side content
    # Vision subtitle
    txbox(slide, Inches(4.7), Inches(1.4), Inches(8.2), Inches(0.6),
          text="연구자 및 R&D 센터를 위한", font_size=Pt(15),
          color=MUTED_COLOR, align=PP_ALIGN.LEFT,
          v_anchor=MSO_ANCHOR.MIDDLE)
    txbox(slide, Inches(4.7), Inches(1.95), Inches(8.2), Inches(0.65),
          text="맞춤형 AI 공동 연구 파트너 플랫폼", font_size=Pt(18), bold=True,
          color=FG_COLOR, align=PP_ALIGN.LEFT,
          v_anchor=MSO_ANCHOR.MIDDLE)

    # Thin separator
    rect(slide, Inches(4.7), Inches(2.75), Inches(7.8), Inches(0.025), TERTIARY_COLOR)

    # Project info block
    txbox_para(slide,
               Inches(4.7), Inches(2.9), Inches(8.0), Inches(1.1),
               [
                   {"text": "비스텔리전스 채용연계형 1기 2차 미니프로젝트",
                    "size": Pt(10.5), "bold": True, "color": FG_COLOR},
                   {"text": "Wrap-Up Report  |  2026. 06. 30",
                    "size": Pt(9.5), "color": MUTED_COLOR, "before": Pt(4)},
               ])

    # Team member cards (3 cards)
    members = [
        ("김지환",  "CS RAG & 문헌 분석"),
        ("신동원",  "Astro RAG & 젬 팩토리"),
        ("천승현",  "Bio RAG & 채팅 허브"),
    ]
    for idx, (name, role) in enumerate(members):
        cx = Inches(4.7) + idx * Inches(2.7)
        cy = Inches(4.3)
        cw = Inches(2.5)
        ch = Inches(1.55)
        rrect(slide, cx, cy, cw, ch)
        card_accent_bar(slide, cx, cy, cw, TERTIARY_COLOR)
        txbox(slide, cx + Inches(0.18), cy + Inches(0.2), cw - Inches(0.36), Inches(0.42),
              text=name, font_size=Pt(13), bold=True, color=FG_COLOR,
              v_anchor=MSO_ANCHOR.MIDDLE)
        txbox(slide, cx + Inches(0.18), cy + Inches(0.68), cw - Inches(0.36), Inches(0.6),
              text=role, font_size=Pt(9.5), color=MUTED_COLOR,
              v_anchor=MSO_ANCHOR.MIDDLE)

    slide_footer(slide, 1)

# =============================================================================
# Slide 2: Project Overview & Milestones
# =============================================================================
def slide_2_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rect(prs, slide, BG_COLOR)
    slide_header(slide, "Project Overview", "프로젝트 개요 및 개발 마일스톤")

    Y0 = Inches(1.55)
    H  = Inches(5.35)

    # ─── Left card ─────────────────────────────────────────────────────────
    LW = Inches(5.5)
    rrect(slide, Inches(0.5), Y0, LW, H)
    card_accent_bar(slide, Inches(0.5), Y0, LW)

    # Vision section
    txbox(slide, Inches(0.75), Y0 + Inches(0.25), LW - Inches(0.5), Inches(0.38),
          text="핵심 비전", font_size=Pt(12), bold=True, color=ACCENT_COLOR,
          v_anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, Inches(0.75), Y0 + Inches(0.65), LW - Inches(0.5), Inches(0.018), TERTIARY_COLOR)

    txbox_para(slide,
               Inches(0.75), Y0 + Inches(0.72), LW - Inches(0.55), Inches(1.7),
               [
                   {"text": "\u201c연구자 및 R&D 센터를 위한 맞춤형 AI 공동 연구 파트너\u201d",
                    "size": Pt(10.5), "bold": True, "color": FG_COLOR, "ls": 1.3},
                   {"text": "파편화된 학술 데이터와 실시간 웹 동향을 단일 지능형 컨텍스트로 통합하여 "
                            "단순 자료 조사 비용을 줄이고 본연의 연구 생산성을 극대화합니다.",
                    "size": Pt(9.5), "color": MUTED_COLOR, "ls": 1.3, "before": Pt(6)},
               ], v_anchor=MSO_ANCHOR.TOP)

    # Team Roles section
    txbox(slide, Inches(0.75), Y0 + Inches(2.55), LW - Inches(0.5), Inches(0.38),
          text="팀 역할 분담", font_size=Pt(12), bold=True, color=ACCENT_COLOR,
          v_anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, Inches(0.75), Y0 + Inches(2.95), LW - Inches(0.5), Inches(0.018), TERTIARY_COLOR)

    roles = [
        ("김지환",  "컴퓨터 과학 분야 RAG 파이프라인 구현 & 대규모 문헌 분석"),
        ("신동원",  "천문학 분야 RAG 파이프라인 구현 & 젬 팩토리 기능 개발"),
        ("천승현",  "생명공학 분야 RAG 파이프라인 구현 & 채팅 허브 기능 개발"),
    ]
    paras = []
    for name, role in roles:
        paras.append({"text": f"• {name}",
                      "size": Pt(10), "bold": True, "color": FG_COLOR, "before": Pt(10)})
        paras.append({"text": f"  {role}",
                      "size": Pt(9.5), "color": MUTED_COLOR, "ls": 1.25, "before": Pt(2)})
    txbox_para(slide, Inches(0.75), Y0 + Inches(3.0), LW - Inches(0.55), Inches(2.1),
               paras, v_anchor=MSO_ANCHOR.TOP)

    # ─── Right card ────────────────────────────────────────────────────────
    RX = Inches(6.3)
    RW = Inches(6.55)
    rrect(slide, RX, Y0, RW, H)
    card_accent_bar(slide, RX, Y0, RW)

    txbox(slide, RX + Inches(0.25), Y0 + Inches(0.25), RW - Inches(0.5), Inches(0.38),
          text="개발 마일스톤  (총 2주)", font_size=Pt(12), bold=True, color=ACCENT_COLOR,
          v_anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, RX + Inches(0.25), Y0 + Inches(0.65), RW - Inches(0.5), Inches(0.018), TERTIARY_COLOR)

    milestones = [
        ("Phase 1",  "6/16 ~ 6/17",  "데이터 인프라 구축",
         "3대 핵심 학술 분야 총 106,974건 대용량 데이터셋 적재 완료"),
        ("Phase 2",  "6/18 ~ 6/26",  "핵심 서비스 구현",
         "병렬 융합 RAG 채팅 허브, 비동기 문헌 분석기, 젬 팩토리 개발 완료"),
        ("Phase 3",  "6/29 ~ 6/30",  "통합 검증 및 배포",
         "전 과정 비동기 및 API 단위 테스트 수행 및 로컬 전체 시스템 배포"),
    ]
    for i, (phase, date, title, desc) in enumerate(milestones):
        my = Y0 + Inches(0.85) + i * Inches(1.48)
        mh = Inches(1.32)
        mw = RW - Inches(0.5)
        # milestone inner card
        mc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    RX + Inches(0.25), my, mw, mh)
        mc.fill.solid(); mc.fill.fore_color.rgb = SIDEBAR_COLOR
        mc.line.color.rgb = BORDER_COLOR; mc.line.width = Pt(1)
        try: mc.adjustments[0] = 0.03
        except: pass

        # Phase badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       RX + Inches(0.38), my + Inches(0.22),
                                       Inches(0.72), Inches(0.30))
        badge.fill.solid(); badge.fill.fore_color.rgb = ACCENT_COLOR
        badge.line.fill.background()
        txbox(slide, RX + Inches(0.38), my + Inches(0.22), Inches(0.72), Inches(0.30),
              text=phase, font_size=Pt(8), bold=True,
              color=WHITE_COLOR, align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

        # Date
        txbox(slide, RX + Inches(1.18), my + Inches(0.22), Inches(2.0), Inches(0.30),
              text=date, font_size=Pt(8.5), color=MUTED_COLOR, v_anchor=MSO_ANCHOR.MIDDLE)

        # Title
        txbox(slide, RX + Inches(0.38), my + Inches(0.57), mw - Inches(0.26), Inches(0.34),
              text=title, font_size=Pt(11.5), bold=True, color=FG_COLOR, v_anchor=MSO_ANCHOR.MIDDLE)

        # Desc
        txbox(slide, RX + Inches(0.38), my + Inches(0.9), mw - Inches(0.26), Inches(0.34),
              text=desc, font_size=Pt(9.5), color=MUTED_COLOR, v_anchor=MSO_ANCHOR.MIDDLE)

    slide_footer(slide, 2)

# =============================================================================
# Slide 3: Lean Canvas
# =============================================================================
def slide_3_leancanvas(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rect(prs, slide, BG_COLOR)
    slide_header(slide, "Business Strategy", "비즈니스 지향점 분석 (Lean Canvas)")

    cells = [
        (Inches(0.5),  Inches(1.55), "문제 & 타겟 고객",
         "• 핵심 문제: 선행 문헌 조사 피로, AI 학술 환각(원문 검증 공수), IP 유출 우려\n"
         "• 타겟 고객: 석·박사 연구원, 대학 교수, R&D 센터 책임연구원, 특허 변리사"),

        (Inches(6.92), Inches(1.55), "솔루션 & 가치 제안",
         "• 솔루션: 듀얼 트랙 병렬 RAG 스트리밍, 연구 공백 분석기, 사용자 정의 Gem 팩토리\n"
         "• 가치 제안: 학술 DB와 실시간 웹 검색의 병렬 융합, 비동기 대규모 문헌 분석 대시보드"),

        (Inches(0.5),  Inches(4.30), "경쟁 우위 & 핵심 지표",
         "• 독점 우위: pgvector HNSW 고속 검색 아키텍처, 100% 팩트 보존 번역 필터 가드\n"
         "• 핵심 지표: 세션당 대화 깊이, 연구 공백 도출 성공률, 인용 출처 CTR"),

        (Inches(6.92), Inches(4.30), "비용 구조 & 수익원",
         "• 비용 구조: LLM API 사용료, pgvector RDS 인프라 비용, 실시간 검색 API 사용료\n"
         "• 수익 구조: Pro 멤버십 ($19.99/월), 기업 전용 보안 On-premise 구축 패키지"),
    ]

    CW = Inches(5.93)
    CH = Inches(2.50)

    for x, y, title, content in cells:
        rrect(slide, x, y, CW, CH)
        card_accent_bar(slide, x, y, CW)

        # Card title
        txbox(slide, x + Inches(0.22), y + Inches(0.12), CW - Inches(0.44), Inches(0.42),
              text=title, font_size=Pt(12), bold=True, color=ACCENT_COLOR,
              v_anchor=MSO_ANCHOR.MIDDLE)

        rect(slide, x + Inches(0.22), y + Inches(0.56), CW - Inches(0.44), Inches(0.018), TERTIARY_COLOR)

        # Content
        txbox(slide, x + Inches(0.22), y + Inches(0.65), CW - Inches(0.44), CH - Inches(0.80),
              text=content, font_size=Pt(9.5), color=FG_COLOR, line_spacing=1.35,
              v_anchor=MSO_ANCHOR.MIDDLE)

    slide_footer(slide, 3)

# =============================================================================
# Slide 4: System Architecture
# =============================================================================
def slide_4_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rect(prs, slide, BG_COLOR)
    slide_header(slide, "System Architecture", "시스템 아키텍처 및 대용량 데이터 관리")

    Y0 = Inches(1.55)
    H  = Inches(5.35)

    # ─── Left: Tech Stack ──────────────────────────────────────────────────
    LW = Inches(5.5)
    rrect(slide, Inches(0.5), Y0, LW, H)
    card_accent_bar(slide, Inches(0.5), Y0, LW)

    txbox(slide, Inches(0.75), Y0 + Inches(0.25), LW - Inches(0.5), Inches(0.38),
          text="3-Tier 물리적 레이어 & 기술 스택",
          font_size=Pt(12), bold=True, color=ACCENT_COLOR, v_anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, Inches(0.75), Y0 + Inches(0.65), LW - Inches(0.5), Inches(0.018), TERTIARY_COLOR)

    techs = [
        ("Frontend",         "Next.js · JavaScript · Axios · Bootstrap 5",
                             "App Router 기반 렌더링, SSE 실시간 통신 지원"),
        ("Backend",          "FastAPI · Python 3.12 · Pydantic v2 DTO",
                             "Async ASGI 기반 스트리밍 및 SSE 브로드캐스팅"),
        ("DB & Vector",      "PostgreSQL 17 · pgvector 0.7+ · SQLAlchemy",
                             "HNSW 인덱싱 기반 3072차원 고차원 벡터 매칭 (68ms)"),
        ("AI Agent",         "LangGraph · LangChain · OpenAI · Tavily",
                             "asyncio.gather 병렬 LLM 오케스트레이션"),
    ]

    for i, (layer, stack, note) in enumerate(techs):
        iy = Y0 + Inches(0.85) + i * Inches(1.1)
        ih = Inches(0.95)
        iw = LW - Inches(0.5)

        ic = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.75), iy, iw, ih)
        ic.fill.solid(); ic.fill.fore_color.rgb = SIDEBAR_COLOR
        ic.line.color.rgb = BORDER_COLOR; ic.line.width = Pt(1)
        try: ic.adjustments[0] = 0.03
        except: pass

        # Left accent
        la = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.75), iy, Inches(0.07), ih)
        la.fill.solid(); la.fill.fore_color.rgb = ACCENT_COLOR
        la.line.fill.background()

        txbox(slide, Inches(0.95), iy + Inches(0.08), iw - Inches(0.28), Inches(0.35),
              text=layer, font_size=Pt(10.5), bold=True, color=FG_COLOR, v_anchor=MSO_ANCHOR.MIDDLE)
        txbox(slide, Inches(0.95), iy + Inches(0.38), iw - Inches(0.28), Inches(0.24),
              text=stack, font_size=Pt(9), bold=True, color=ACCENT_COLOR, v_anchor=MSO_ANCHOR.MIDDLE)
        txbox(slide, Inches(0.95), iy + Inches(0.60), iw - Inches(0.28), Inches(0.28),
              text=note, font_size=Pt(8.5), color=MUTED_COLOR, v_anchor=MSO_ANCHOR.MIDDLE)

    # ─── Right: Data ───────────────────────────────────────────────────────
    RX = Inches(6.3)
    RW = Inches(6.55)
    rrect(slide, RX, Y0, RW, H)
    card_accent_bar(slide, RX, Y0, RW)

    txbox(slide, RX + Inches(0.25), Y0 + Inches(0.25), RW - Inches(0.5), Inches(0.38),
          text="대용량 ArXiv 학술 데이터 적재 현황",
          font_size=Pt(12), bold=True, color=ACCENT_COLOR, v_anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, RX + Inches(0.25), Y0 + Inches(0.65), RW - Inches(0.5), Inches(0.018), TERTIARY_COLOR)

    # Domain stats: 3 cards
    domains = [
        ("생명과학  q-bio.GN",   "54,066건", "가장 큰 임베딩 공간, 유전자 편집 및 단백질 분석 최적화"),
        ("천문학  astro-ph.EP", "35,083건", "외계행성 및 궤도역학 이론 검증용 대형 임베딩 공간"),
        ("컴퓨터공학  cs.NE",   "17,825건", "신경망 진화 및 딥러닝 핵심 개념 탐색 고속 엔진"),
    ]
    for i, (domain, count, note) in enumerate(domains):
        dy = Y0 + Inches(0.85) + i * Inches(1.1)
        dh = Inches(0.95)
        dw = RW - Inches(0.5)
        dc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    RX + Inches(0.25), dy, dw, dh)
        dc.fill.solid(); dc.fill.fore_color.rgb = SIDEBAR_COLOR
        dc.line.color.rgb = BORDER_COLOR; dc.line.width = Pt(1)
        try: dc.adjustments[0] = 0.03
        except: pass

        la = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    RX + Inches(0.25), dy, Inches(0.07), dh)
        la.fill.solid(); la.fill.fore_color.rgb = TERTIARY_COLOR
        la.line.fill.background()

        txbox(slide, RX + Inches(0.45), dy + Inches(0.08), dw - Inches(2.5), Inches(0.35),
              text=domain, font_size=Pt(10.5), bold=True, color=FG_COLOR, v_anchor=MSO_ANCHOR.MIDDLE)
        txbox(slide, RX + Inches(0.45) + dw - Inches(2.5), dy + Inches(0.08), Inches(2.2), Inches(0.35),
              text=count, font_size=Pt(14), bold=True, color=ACCENT_COLOR,
              align=PP_ALIGN.RIGHT, v_anchor=MSO_ANCHOR.MIDDLE)
        txbox(slide, RX + Inches(0.45), dy + Inches(0.54), dw - Inches(0.28), Inches(0.36),
              text=note, font_size=Pt(9), color=MUTED_COLOR, v_anchor=MSO_ANCHOR.MIDDLE)

    # Embedding spec banner
    by = Y0 + Inches(4.2)
    bw = RW - Inches(0.5)
    bb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                RX + Inches(0.25), by, bw, Inches(0.8))
    bb.fill.solid(); bb.fill.fore_color.rgb = ACCENT_COLOR
    bb.line.fill.background()
    txbox(slide, RX + Inches(0.4), by, bw - Inches(0.3), Inches(0.8),
          text="✦  OpenAI text-embedding-3-large  ·  3,072차원  ·  HNSW cosine  ·  avg. 68ms",
          font_size=Pt(9.5), bold=True, color=WHITE_COLOR,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

    slide_footer(slide, 4)

# =============================================================================
# Slide 5: Core Features — Chat Hub & Gap Analyzer
# =============================================================================
def slide_5_features1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rect(prs, slide, BG_COLOR)
    slide_header(slide, "Core Features  01 & 02", "병렬 융합 Q&A 채팅 허브 및 연구 공백 분석기")

    Y0 = Inches(1.55)
    H  = Inches(5.35)
    CW = Inches(6.0)

    for col, (feat_num, feat_title, intro, bullets, badge_color) in enumerate([
        ("01", "일반 채팅 허브 (Q&A Hub)",
         "자연어 질문에 대해 논문 RAG와 실시간 웹 검색을 병렬로 구동하여 두 지식을 최적으로 융합한 답변을 실시간 스트리밍합니다.",
         [
             ("무조건적 병렬 인출",
              "asyncio.gather를 이용하여 DB 검색과 웹 검색을 동시 처리 — 순차 대비 평균 23% 대기 지연 감소"),
             ("SSE 실시간 스트리밍",
              "HTTP SSE 방식으로 토큰 발생 시마다 실시간 푸시, 타자 효과 마크다운 렌더링"),
             ("인용 출처 매핑 & 적재",
              "참조된 ArXiv 논문 서지 메타데이터를 1:1 인덱싱하여 chat_source 테이블에 영구 적재"),
         ], ACCENT_COLOR),

        ("02", "연구 공백 분석기 (Gap Analyzer)",
         "수십 편의 선행 연구 데이터를 배치로 분석하여, 각 문헌의 한계점을 취합하고 미개척 연구 가설을 자동 도출합니다.",
         [
             ("비동기 BackgroundTasks",
              "task_id 즉시 발급 후 연산을 BackgroundTasks로 오프로딩, 진행률(10→40→80→100%) 실시간 갱신"),
             ("100% 팩트 보존 번역 가드",
              "한글 번역 과정에서 핵심 영문 인용구를 파이썬 레이어에서 백업 후 번역 완료 시 오버라이트 복원"),
             ("SSE 완료 알림 & Redis 캐싱",
              "분석 완료 즉시 SSE 브로드캐스트로 토스트 알림 노출 및 번역 결과 JSONB 캐시 영구 저장"),
         ], DEEP_ACCENT),
    ]):
        x = Inches(0.5) + col * Inches(6.42)
        rrect(slide, x, Y0, CW, H)
        card_accent_bar(slide, x, Y0, CW, badge_color)

        # Feature number badge
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x + Inches(0.22), Y0 + Inches(0.2), Inches(0.5), Inches(0.3))
        b.fill.solid(); b.fill.fore_color.rgb = badge_color; b.line.fill.background()
        txbox(slide, x + Inches(0.22), Y0 + Inches(0.2), Inches(0.5), Inches(0.3),
              text=feat_num, font_size=Pt(9), bold=True, color=WHITE_COLOR,
              align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

        txbox(slide, x + Inches(0.82), Y0 + Inches(0.2), CW - Inches(1.1), Inches(0.3),
              text=feat_title, font_size=Pt(12), bold=True, color=FG_COLOR,
              v_anchor=MSO_ANCHOR.MIDDLE)

        rect(slide, x + Inches(0.22), Y0 + Inches(0.60), CW - Inches(0.44), Inches(0.018), TERTIARY_COLOR)

        txbox(slide, x + Inches(0.22), Y0 + Inches(0.68), CW - Inches(0.44), Inches(0.85),
              text=intro, font_size=Pt(9.5), color=MUTED_COLOR, line_spacing=1.3,
              v_anchor=MSO_ANCHOR.TOP)

        for j, (bul_title, bul_desc) in enumerate(bullets):
            by = Y0 + Inches(1.6) + j * Inches(1.18)
            bh = Inches(1.05)
            bw = CW - Inches(0.44)
            bc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x + Inches(0.22), by, bw, bh)
            bc.fill.solid(); bc.fill.fore_color.rgb = SIDEBAR_COLOR
            bc.line.color.rgb = BORDER_COLOR; bc.line.width = Pt(1)
            try: bc.adjustments[0] = 0.03
            except: pass

            # tiny left bar
            lb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        x + Inches(0.22), by, Inches(0.06), bh)
            lb.fill.solid(); lb.fill.fore_color.rgb = badge_color; lb.line.fill.background()

            txbox(slide, x + Inches(0.38), by + Inches(0.08), bw - Inches(0.2), Inches(0.34),
                  text=bul_title, font_size=Pt(10), bold=True, color=FG_COLOR,
                  v_anchor=MSO_ANCHOR.MIDDLE)
            txbox(slide, x + Inches(0.38), by + Inches(0.44), bw - Inches(0.2), Inches(0.52),
                  text=bul_desc, font_size=Pt(9), color=MUTED_COLOR, line_spacing=1.25,
                  v_anchor=MSO_ANCHOR.MIDDLE)

    slide_footer(slide, 5)

# =============================================================================
# Slide 6: Gem Factory (4-step flow)
# =============================================================================
def slide_6_gemfactory(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rect(prs, slide, BG_COLOR)
    slide_header(slide, "Core Features  03", "맞춤형 AI 연구 비서 — Gem 팩토리 파이프라인")

    # Intro bar
    intro = ("사용자가 연구 카테고리와 시스템 프롬프트(페르소나)를 조합하여 전용 연구 비서(Gem)를 개설하고, "
             "독립된 격리 데이터 환경에서 RAG 대화를 나눌 수 있는 4단계 파이프라인입니다.")
    ib = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(1.5), Inches(12.33), Inches(0.68))
    ib.fill.solid(); ib.fill.fore_color.rgb = TERTIARY_COLOR; ib.line.fill.background()
    txbox(slide, Inches(0.65), Inches(1.5), Inches(12.0), Inches(0.68),
          text=intro, font_size=Pt(10), color=FG_COLOR,
          v_anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)

    steps = [
        ("Step 1", "페르소나 바인딩",
         "사용자가 설정한 RAG 참고 분야(CS / Bio / Astro) 및 연구 지침 시스템 프롬프트를 결합해 "
         "젬의 고유 메타데이터 레코드를 생성하고 gem 테이블에 영구 등록합니다."),
        ("Step 2", "PDF 테넌시 격리",
         "개인 연구 PDF를 업로드하면 백그라운드에서 500자 청킹 후 text-embedding-3-large로 벡터화, "
         "gem_{gem_id}_files 전용 컬렉션에 적재하여 타 사용자와 완벽히 물리 격리합니다."),
        ("Step 3", "동적 도구 주입",
         "세션 개설 시 렉시컬 스코프 내에 gem_id를 가두는 클로저(Closure) 함수를 동적으로 빌드하여 "
         "Supervisor Agent의 툴로 실시간 주입, 지정 컬렉션에서만 검색하도록 강제 제어합니다."),
        ("Step 4", "물리 완전 소멸",
         "젬 삭제 시 DB Cascade 연쇄 삭제 제약에 의해 메타데이터가 드롭되는 동시에, "
         "gem_{gem_id}_files 테이블을 'DROP TABLE' 처리하여 디스크에서 0바이트로 완전 소각합니다."),
    ]

    CW = Inches(2.75)
    CH = Inches(4.55)
    gap = Inches(0.3)
    x0 = Inches(0.5)
    y0 = Inches(2.35)

    for i, (step, title, desc) in enumerate(steps):
        x = x0 + i * (CW + gap)
        rrect(slide, x, y0, CW, CH)
        card_accent_bar(slide, x, y0, CW)

        # Step badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x + Inches(0.2), y0 + Inches(0.2), Inches(0.85), Inches(0.3))
        badge.fill.solid(); badge.fill.fore_color.rgb = ACCENT_COLOR; badge.line.fill.background()
        txbox(slide, x + Inches(0.2), y0 + Inches(0.2), Inches(0.85), Inches(0.3),
              text=step, font_size=Pt(8.5), bold=True, color=WHITE_COLOR,
              align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

        txbox(slide, x + Inches(0.2), y0 + Inches(0.62), CW - Inches(0.4), Inches(0.46),
              text=title, font_size=Pt(13), bold=True, color=FG_COLOR,
              v_anchor=MSO_ANCHOR.MIDDLE)

        rect(slide, x + Inches(0.2), y0 + Inches(1.12), CW - Inches(0.4), Inches(0.02), TERTIARY_COLOR)

        txbox(slide, x + Inches(0.2), y0 + Inches(1.2), CW - Inches(0.4), CH - Inches(1.4),
              text=desc, font_size=Pt(9), color=MUTED_COLOR, line_spacing=1.35,
              v_anchor=MSO_ANCHOR.TOP)

        # Arrow connector (except last)
        if i < len(steps) - 1:
            ax = x + CW + Inches(0.05)
            ay = y0 + CH / 2 - Inches(0.12)
            txbox(slide, ax, ay, gap - Inches(0.05), Inches(0.25),
                  text="▶", font_size=Pt(11), color=ACCENT_COLOR,
                  align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

    slide_footer(slide, 6)

# =============================================================================
# Slide 7: Troubleshooting (3 Cases)
# =============================================================================
def slide_7_troubleshooting(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rect(prs, slide, BG_COLOR)
    slide_header(slide, "Troubleshooting", "핵심 기술 문제 해결 및 리팩토링 사례")

    cases = [
        ("Case 01", "Structured Outputs &\nStreamingResponse 결합",
         "이벤트 기반 NDJSON 스트리밍 도입",
         "문제",
         "LLM 구조화 모델은 완결된 JSON이 완성되기 전까지 토큰 단위 실시간 스트리밍이 불가능하며, 두 방식을 무리하게 혼용할 경우 타입 불일치 및 클라이언트 파싱 오류가 발생함.",
         "해결",
         "이벤트 기반 NDJSON 스트리밍을 구현하여 답변 본문은 즉시 토큰 단위로 전송하고, 추천 질문 및 RAG 출처 등 구조화 정보는 스트림 완료 시점에 Pydantic DTO 검증 후 최종 이벤트 라인으로 병합 전송하여 극복했습니다."),

        ("Case 02", "10만 건 대용량 적재\n커넥션 락업",
         "페이징 오프셋 분할 커밋 적용",
         "문제",
         "ArXiv 10만 건 이상의 벡터 임베딩 데이터를 단일 DB 트랜잭션으로 적재할 때, 메모리 누수와 데이터베이스 커넥션 락(Connection Lock)이 걸려 타임아웃 오류가 빈번히 일어남.",
         "해결",
         "5,000건 단위로 트랜잭션을 분할 및 배치 Commit하고 주기적으로 메모리를 GC 처리함으로써 106,974건의 데이터를 단 한 번의 다운타임 없이 안정적으로 적재하는 데 성공했습니다."),

        ("Case 03", "번역 과정의 학술\n인용구 원문 왜곡",
         "Python 레이어 원문 백업/복원 가드",
         "문제",
         "GPT-4o 번역 단계에서 학술적으로 엄격해야 하는 원어 인용구(source_quote)들이 자연스러운 한글로 오역되거나 원어 고유 코드가 손실되어, 원문 출처를 역추적하는 신뢰성에 손상이 감.",
         "해결",
         "번역 전 파이썬 코드 레이어에서 핵심 인용구 배열을 복사(Backup)하고, 번역 완성 후 원어를 강제 주입(Restore)하는 하이브리드 필터를 설계하여 인용구 유실률 0%를 달성했습니다."),
    ]

    CW = Inches(3.78)
    CH = Inches(5.35)
    x0 = Inches(0.5)
    y0 = Inches(1.55)
    gap = Inches(0.26)

    for i, (case_no, title, subtitle, prob_label, prob, sol_label, sol) in enumerate(cases):
        x = x0 + i * (CW + gap)
        rrect(slide, x, y0, CW, CH)
        card_accent_bar(slide, x, y0, CW)

        # Case badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x + Inches(0.22), y0 + Inches(0.2), Inches(0.82), Inches(0.3))
        badge.fill.solid(); badge.fill.fore_color.rgb = ACCENT_COLOR; badge.line.fill.background()
        txbox(slide, x + Inches(0.22), y0 + Inches(0.2), Inches(0.82), Inches(0.3),
              text=case_no, font_size=Pt(8.5), bold=True, color=WHITE_COLOR,
              align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

        # Title
        txbox(slide, x + Inches(0.22), y0 + Inches(0.58), CW - Inches(0.44), Inches(0.7),
              text=title, font_size=Pt(12), bold=True, color=FG_COLOR, line_spacing=1.25,
              v_anchor=MSO_ANCHOR.MIDDLE)

        # Subtitle (approach tag)
        txbox(slide, x + Inches(0.22), y0 + Inches(1.32), CW - Inches(0.44), Inches(0.3),
              text=f"→ {subtitle}", font_size=Pt(9.5), bold=True, color=ACCENT_COLOR,
              v_anchor=MSO_ANCHOR.MIDDLE)

        rect(slide, x + Inches(0.22), y0 + Inches(1.68), CW - Inches(0.44), Inches(0.018), TERTIARY_COLOR)

        # Problem block
        pb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x + Inches(0.22), y0 + Inches(1.8), CW - Inches(0.44), Inches(1.3))
        pb.fill.solid(); pb.fill.fore_color.rgb = SIDEBAR_COLOR
        pb.line.color.rgb = BORDER_COLOR; pb.line.width = Pt(1)
        try: pb.adjustments[0] = 0.03
        except: pass

        pl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x + Inches(0.22), y0 + Inches(1.8), Inches(0.06), Inches(1.3))
        pl.fill.solid(); pl.fill.fore_color.rgb = RGBColor(220, 100, 80); pl.line.fill.background()

        txbox(slide, x + Inches(0.38), y0 + Inches(1.83), CW - Inches(0.66), Inches(0.28),
              text="⚑  " + prob_label, font_size=Pt(8.5), bold=True,
              color=RGBColor(200, 80, 60), v_anchor=MSO_ANCHOR.MIDDLE)
        txbox(slide, x + Inches(0.38), y0 + Inches(2.1), CW - Inches(0.66), Inches(0.95),
              text=prob, font_size=Pt(8.5), color=MUTED_COLOR, line_spacing=1.25,
              v_anchor=MSO_ANCHOR.TOP)

        # Solution block
        sb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x + Inches(0.22), y0 + Inches(3.22), CW - Inches(0.44), Inches(1.84))
        sb.fill.solid(); sb.fill.fore_color.rgb = SIDEBAR_COLOR
        sb.line.color.rgb = BORDER_COLOR; sb.line.width = Pt(1)
        try: sb.adjustments[0] = 0.03
        except: pass

        sl_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x + Inches(0.22), y0 + Inches(3.22), Inches(0.06), Inches(1.84))
        sl_bar.fill.solid(); sl_bar.fill.fore_color.rgb = ACCENT_COLOR; sl_bar.line.fill.background()

        txbox(slide, x + Inches(0.38), y0 + Inches(3.26), CW - Inches(0.66), Inches(0.28),
              text="✔  " + sol_label, font_size=Pt(8.5), bold=True,
              color=ACCENT_COLOR, v_anchor=MSO_ANCHOR.MIDDLE)
        txbox(slide, x + Inches(0.38), y0 + Inches(3.54), CW - Inches(0.66), Inches(1.42),
              text=sol, font_size=Pt(8.5), color=MUTED_COLOR, line_spacing=1.25,
              v_anchor=MSO_ANCHOR.TOP)

    slide_footer(slide, 7)

# =============================================================================
# Slide 8: Future Roadmap
# =============================================================================
def slide_8_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_rect(prs, slide, BG_COLOR)
    slide_header(slide, "Future Roadmap", "한계 극복 및 향후 발전 로드맵")

    roadmaps = [
        ("Roadmap 01", "기업용 보안 샌드박스",
         "물리 파쇄 (Shredding) 데몬 완성",
         [
             ("한계점",
              "보안 비서 젬 및 업로드된 PDF 파일이 런타임 ID 분리에만 의존하여, 무활동 상태에서도 디스크 및 DB에 데이터 잔재가 남아 정보 유출 취약점이 존재함.",
              False),
             ("보완 계획",
              "• Redis TTL 만료 이벤트 감지 → 비동기 삭제 트리거\n"
              "• MacOS shred 커맨드로 난수 3회 덮어쓰기 파괴\n"
              "• ON DELETE CASCADE로 pgvector 임시 테이블 연쇄 파쇄",
              True),
         ]),
        ("Roadmap 02", "다중 에이전트 교차 검증",
         "LangGraph 피어 리뷰 토론 엔진",
         [
             ("한계점",
              "현재 RAG 합성의 최종 조율이 단일 Synthesis 노드에서만 이루어지기 때문에 잠재적 할루시네이션(환각) 리스크를 완전히 상쇄할 수 없음.",
              False),
             ("보완 계획",
              "• 방법론 검증자: 가설의 논리적 모순 공격\n"
              "• 신규성 분석자: 특허 RAG 대조를 통한 선행 중복도 검사\n"
              "• LangGraph 순환형 그래프 — 2회 이상 피드백 패스 후 최종 합의 리포트 발행",
              True),
         ]),
        ("Roadmap 03", "특허 및 이종 산업 융합",
         "BM25 + HNSW 하이브리드 RRF",
         [
             ("한계점",
              "현재 생명공학, 천문학, 컴퓨터 과학 3대 분야의 ArXiv DB로 타겟이 한정되어 있어, 특허 변리 및 타 산업군 R&D로의 실질적 도메인 확장이 필요함.",
              False),
             ("보완 계획",
              "• BM25 키워드 검색 + HNSW 시맨틱 검색 → RRF 상호 순위 융합 노드\n"
              "• Dynamic JSONB DB 스키마 로더로 무중단 도메인 유입 지원\n"
              "• 법령, 특허, 신소재 등 전문 지식 레이어 순차 통합",
              True),
         ]),
    ]

    CW = Inches(3.78)
    CH = Inches(5.35)
    x0 = Inches(0.5)
    y0 = Inches(1.55)
    gap = Inches(0.26)

    for i, (road_no, title, subtitle, blocks) in enumerate(roadmaps):
        x = x0 + i * (CW + gap)
        rrect(slide, x, y0, CW, CH)
        card_accent_bar(slide, x, y0, CW)

        # Roadmap badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x + Inches(0.22), y0 + Inches(0.2), Inches(1.05), Inches(0.3))
        badge.fill.solid(); badge.fill.fore_color.rgb = ACCENT_COLOR; badge.line.fill.background()
        txbox(slide, x + Inches(0.22), y0 + Inches(0.2), Inches(1.05), Inches(0.3),
              text=road_no, font_size=Pt(8.5), bold=True, color=WHITE_COLOR,
              align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

        txbox(slide, x + Inches(0.22), y0 + Inches(0.58), CW - Inches(0.44), Inches(0.44),
              text=title, font_size=Pt(13), bold=True, color=FG_COLOR,
              v_anchor=MSO_ANCHOR.MIDDLE)

        txbox(slide, x + Inches(0.22), y0 + Inches(1.04), CW - Inches(0.44), Inches(0.3),
              text=f"→ {subtitle}", font_size=Pt(9.5), bold=True, color=ACCENT_COLOR,
              v_anchor=MSO_ANCHOR.MIDDLE)

        rect(slide, x + Inches(0.22), y0 + Inches(1.42), CW - Inches(0.44), Inches(0.018), TERTIARY_COLOR)

        block_heights = [Inches(1.35), Inches(2.5)]
        block_y = y0 + Inches(1.55)
        for j, (blabel, btext, is_sol) in enumerate(blocks):
            bh = block_heights[j]
            bw = CW - Inches(0.44)
            bc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x + Inches(0.22), block_y, bw, bh)
            bc.fill.solid(); bc.fill.fore_color.rgb = SIDEBAR_COLOR
            bc.line.color.rgb = BORDER_COLOR; bc.line.width = Pt(1)
            try: bc.adjustments[0] = 0.03
            except: pass

            bar_color = ACCENT_COLOR if is_sol else RGBColor(220, 100, 80)
            bl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x + Inches(0.22), block_y, Inches(0.06), bh)
            bl.fill.solid(); bl.fill.fore_color.rgb = bar_color; bl.line.fill.background()

            label_color = ACCENT_COLOR if is_sol else RGBColor(200, 80, 60)
            label_icon = "✔" if is_sol else "⚑"
            txbox(slide, x + Inches(0.38), block_y + Inches(0.1), bw - Inches(0.2), Inches(0.28),
                  text=f"{label_icon}  {blabel}", font_size=Pt(8.5), bold=True,
                  color=label_color, v_anchor=MSO_ANCHOR.MIDDLE)
            txbox(slide, x + Inches(0.38), block_y + Inches(0.38), bw - Inches(0.2), bh - Inches(0.52),
                  text=btext, font_size=Pt(8.5), color=MUTED_COLOR, line_spacing=1.3,
                  v_anchor=MSO_ANCHOR.TOP)

            block_y += bh + Inches(0.12)

    slide_footer(slide, 8)

# =============================================================================
# Main
# =============================================================================
def main():
    prs = init_presentation()
    slide_1_cover(prs)
    slide_2_overview(prs)
    slide_3_leancanvas(prs)
    slide_4_architecture(prs)
    slide_5_features1(prs)
    slide_6_gemfactory(prs)
    slide_7_troubleshooting(prs)
    slide_8_roadmap(prs)

    out_dir  = "/Users/pileuszu/Repos/bist-mini-2/docs/deliverables/final"
    out_path = os.path.join(out_dir, "project_presentation.pptx")
    os.makedirs(out_dir, exist_ok=True)
    prs.save(out_path)
    print(f"✅  Saved: {out_path}")

if __name__ == "__main__":
    main()
